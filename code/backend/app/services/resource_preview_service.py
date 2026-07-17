"""Safe, bounded renderers for files shown in the student resource library."""

from __future__ import annotations

import html
import zipfile
from dataclasses import dataclass
from pathlib import Path

from docx import Document
from pptx import Presentation


class ResourcePreviewError(ValueError):
    """A file could not be prepared for a safe in-app preview."""


@dataclass(frozen=True)
class ResourcePreview:
    mode: str
    media_type: str
    content: str | None = None
    stream_file: bool = False


class ResourcePreviewService:
    """Render only the small, inert subset needed by the preview dialog."""

    _MAX_FILE_BYTES = 25 * 1024 * 1024
    _MAX_ARCHIVE_FILES = 800
    _MAX_ARCHIVE_UNCOMPRESSED_BYTES = 40 * 1024 * 1024
    _MAX_TEXT_CHARS = 240_000
    _MAX_PARAGRAPHS = 1_200
    _MAX_TABLE_ROWS = 120
    _MAX_TABLE_COLUMNS = 24
    _MAX_SLIDES = 100

    def prepare(self, path: Path) -> ResourcePreview:
        if not path.is_file():
            raise ResourcePreviewError("资源文件不存在或已被删除")
        if path.stat().st_size > self._MAX_FILE_BYTES:
            raise ResourcePreviewError("文件过大，暂不支持在线预览，请下载后查看")

        extension = path.suffix.lower()
        if extension == ".pdf":
            return ResourcePreview("pdf", "application/pdf", stream_file=True)
        if extension in {".png", ".jpg", ".jpeg"}:
            media_type = "image/png" if extension == ".png" else "image/jpeg"
            return ResourcePreview("image", media_type, stream_file=True)
        if extension == ".docx":
            self._validate_office_archive(path)
            return ResourcePreview("document", "text/html; charset=utf-8", self._document_html(path))
        if extension == ".pptx":
            self._validate_office_archive(path)
            return ResourcePreview("presentation", "text/html; charset=utf-8", self._presentation_html(path))
        if extension in {".md", ".markdown", ".txt"}:
            return ResourcePreview("text", "text/html; charset=utf-8", self._text_html(path))
        # Legacy binary office formats are deliberately not handed to a local
        # converter or shell command.  That is both safer and more predictable.
        raise ResourcePreviewError("该文件格式暂不支持在线预览，请下载后查看")

    def _validate_office_archive(self, path: Path) -> None:
        try:
            with zipfile.ZipFile(path) as archive:
                entries = archive.infolist()
                total = sum(entry.file_size for entry in entries)
        except (OSError, zipfile.BadZipFile) as exc:
            raise ResourcePreviewError("文件无法安全解析，请下载后查看") from exc
        if len(entries) > self._MAX_ARCHIVE_FILES or total > self._MAX_ARCHIVE_UNCOMPRESSED_BYTES:
            raise ResourcePreviewError("文件内容过大，暂不支持在线预览，请下载后查看")

    def _document_html(self, path: Path) -> str:
        try:
            document = Document(str(path))
        except Exception as exc:
            raise ResourcePreviewError("Word 文件无法安全解析，请下载后查看") from exc

        parts: list[str] = []
        active_list: str | None = None
        remaining = self._MAX_TEXT_CHARS
        for paragraph in document.paragraphs[: self._MAX_PARAGRAPHS]:
            text, remaining = self._bounded_text(paragraph.text, remaining)
            if not text:
                continue
            style = (getattr(paragraph.style, "name", "") or "").casefold()
            if "list" in style or "列表" in style:
                list_tag = "ol" if "number" in style or "编号" in style else "ul"
                if active_list != list_tag:
                    if active_list:
                        parts.append(f"</{active_list}>")
                    parts.append(f"<{list_tag}>")
                    active_list = list_tag
                parts.append(f"<li>{html.escape(text)}</li>")
                continue
            if active_list:
                parts.append(f"</{active_list}>")
                active_list = None
            if style.startswith("heading") or style.startswith("标题"):
                level = next((char for char in style if char.isdigit()), "2")
                parts.append(f"<h{min(max(int(level), 1), 4)}>{html.escape(text)}</h{min(max(int(level), 1), 4)}>")
            else:
                parts.append(f"<p>{html.escape(text)}</p>")
            if remaining <= 0:
                break
        if active_list:
            parts.append(f"</{active_list}>")

        for table in document.tables:
            if remaining <= 0:
                break
            rows: list[str] = []
            for row_index, row in enumerate(table.rows[: self._MAX_TABLE_ROWS]):
                cells: list[str] = []
                for cell in row.cells[: self._MAX_TABLE_COLUMNS]:
                    text, remaining = self._bounded_text(cell.text, remaining)
                    tag = "th" if row_index == 0 else "td"
                    cells.append(f"<{tag}>{html.escape(text)}</{tag}>")
                rows.append(f"<tr>{''.join(cells)}</tr>")
            if rows:
                parts.append(f"<table><tbody>{''.join(rows)}</tbody></table>")
        return self._html_shell("Word 文档预览", "document", "".join(parts) or "<p>文档没有可显示的文字内容。</p>")

    def _presentation_html(self, path: Path) -> str:
        try:
            presentation = Presentation(str(path))
        except Exception as exc:
            raise ResourcePreviewError("PPTX 文件无法安全解析，请下载后查看") from exc
        slides: list[str] = []
        remaining = self._MAX_TEXT_CHARS
        for index, slide in enumerate(presentation.slides[: self._MAX_SLIDES], start=1):
            blocks: list[str] = []
            for shape in slide.shapes:
                if remaining <= 0:
                    break
                if getattr(shape, "has_table", False):
                    rows: list[str] = []
                    for row_index, row in enumerate(shape.table.rows[: self._MAX_TABLE_ROWS]):
                        cells: list[str] = []
                        for cell in row.cells[: self._MAX_TABLE_COLUMNS]:
                            text, remaining = self._bounded_text(cell.text, remaining)
                            tag = "th" if row_index == 0 else "td"
                            cells.append(f"<{tag}>{html.escape(text)}</{tag}>")
                        rows.append(f"<tr>{''.join(cells)}</tr>")
                    blocks.append(f"<table><tbody>{''.join(rows)}</tbody></table>")
                elif getattr(shape, "has_text_frame", False):
                    for paragraph in shape.text_frame.paragraphs:
                        text, remaining = self._bounded_text(paragraph.text, remaining)
                        if text:
                            blocks.append(f"<p>{html.escape(text)}</p>")
            slides.append(
                f'<section class="slide"><header>第 {index} 页</header>'
                f"<div>{''.join(blocks) or '<p>此页没有可显示的文字内容。</p>'}</div></section>"
            )
            if remaining <= 0:
                break
        return self._html_shell("PPTX 演示文稿预览", "presentation", "".join(slides))

    def _text_html(self, path: Path) -> str:
        raw = b""
        with path.open("rb") as file:
            raw = file.read(self._MAX_TEXT_CHARS * 4 + 1)
        for encoding in ("utf-8-sig", "utf-8", "gb18030"):
            try:
                text = raw.decode(encoding)
                break
            except UnicodeDecodeError:
                continue
        else:
            text = raw.decode("utf-8", errors="replace")
        text, _ = self._bounded_text(text, self._MAX_TEXT_CHARS)
        return self._html_shell("文本资料预览", "text", f"<pre>{html.escape(text)}</pre>")

    @staticmethod
    def _bounded_text(value: str, remaining: int) -> tuple[str, int]:
        compact = value.strip()
        if remaining <= 0:
            return "", 0
        if len(compact) > remaining:
            return compact[:remaining] + "…", 0
        return compact, remaining - len(compact)

    @staticmethod
    def _html_shell(title: str, mode: str, body: str) -> str:
        escaped_title = html.escape(title)
        return f"""<!doctype html>
<html lang="zh-CN"><head><meta charset="utf-8"><meta name="viewport" content="width=device-width, initial-scale=1">
<meta http-equiv="Content-Security-Policy" content="default-src 'none'; style-src 'unsafe-inline'; img-src data:; base-uri 'none'; form-action 'none'">
<title>{escaped_title}</title><style>
*,*::before,*::after{{box-sizing:border-box}}html,body{{min-width:0}}body{{margin:0;padding:32px;background:#f5f6f8;color:#202939;font:16px/1.75 system-ui,-apple-system,"PingFang SC","Microsoft YaHei",sans-serif}}main{{width:min(210mm,100%);min-height:297mm;margin:auto;padding:22mm 19mm;background:#fff;box-shadow:0 1px 3px #10182812;overflow-wrap:anywhere}}h1,h2,h3,h4{{line-height:1.35;color:#101828}}h1{{font-size:1.8em}}h2{{font-size:1.45em;margin-top:1.8em}}p{{margin:0 0 1em}}ul,ol{{padding-left:1.6em}}table{{width:100%;margin:1.25em 0;border-collapse:collapse;font-size:.92em}}th,td{{padding:8px 10px;border:1px solid #d9dee8;text-align:left;vertical-align:top;white-space:pre-wrap;word-break:break-word}}th{{background:#f4f6fa}}pre{{margin:0;white-space:pre-wrap;overflow-wrap:anywhere;font:14px/1.7 ui-monospace,SFMono-Regular,Menlo,monospace}}.preview-presentation{{width:min(1200px,100%);min-height:0;padding:28px;background:transparent;box-shadow:none}}.slide{{aspect-ratio:16/9;margin:0 auto 28px;padding:5.5% 7%;overflow:auto;border:1px solid #d9dee8;background:#fff;box-shadow:0 8px 24px #10182812}}.slide>header{{margin-bottom:1.2em;color:#667085;font-size:.8em;font-weight:700}}@media(max-width:720px){{body{{padding:12px}}main{{width:100%;min-height:0;padding:24px 20px}}.preview-presentation{{padding:0}}}}@media print{{body{{padding:0;background:#fff}}main{{width:100%;min-height:0;padding:0;box-shadow:none}}.slide{{break-after:page;margin:0;border:0;box-shadow:none}}}}
</style></head><body><main class="preview-{mode}">{body}</main></body></html>"""


resource_preview_service = ResourcePreviewService()
