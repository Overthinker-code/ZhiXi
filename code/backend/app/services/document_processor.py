from typing import List, Optional
from pathlib import Path

import docx
import PyPDF2
try:
    import pdfplumber
except Exception:  # pragma: no cover - optional dependency
    pdfplumber = None
from pptx import Presentation

from app.core.config import settings

# 尝试导入 langchain，失败时提供备用方案
try:
    from langchain_text_splitters import RecursiveCharacterTextSplitter
    LANGCHAIN_AVAILABLE = True
except Exception:
    LANGCHAIN_AVAILABLE = False
    RecursiveCharacterTextSplitter = None

from langchain_core.documents import Document


class DocumentProcessor:
    def __init__(self):
        if LANGCHAIN_AVAILABLE and RecursiveCharacterTextSplitter:
            self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=settings.RAG_CHUNK_SIZE,
            chunk_overlap=settings.RAG_CHUNK_OVERLAP,
            length_function=len,
        )

    def process_word(self, file_path: str) -> List[Document]:
        doc = docx.Document(file_path)
        text = "\n".join([paragraph.text for paragraph in doc.paragraphs])
        metadata = self._build_metadata(file_path, doc_type="word")
        return self._split_text(text, metadata=metadata)

    def process_pdf(self, file_path: str) -> List[Document]:
        text = ""
        if pdfplumber is not None:
            with pdfplumber.open(file_path) as pdf:
                for page in pdf.pages:
                    text += page.extract_text() or ""
        else:
            with open(file_path, "rb") as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page in pdf_reader.pages:
                    text += page.extract_text() or ""
        metadata = self._build_metadata(file_path, doc_type="pdf")
        return self._split_text(text, metadata=metadata)

    def process_ppt(self, file_path: str) -> List[Document]:
        prs = Presentation(file_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        metadata = self._build_metadata(file_path, doc_type="ppt")
        return self._split_text(text, metadata=metadata)

    def process_markdown(self, file_path: str) -> List[Document]:
        text = self._read_plain_text(file_path)
        metadata = self._build_metadata(file_path, doc_type="markdown")
        return self._split_text(text, metadata=metadata)

    def _read_plain_text(self, file_path: str) -> str:
        for encoding in ("utf-8-sig", "utf-8", "gb18030"):
            try:
                with open(file_path, "r", encoding=encoding) as file:
                    return file.read()
            except UnicodeDecodeError:
                continue
        with open(file_path, "r", encoding="utf-8", errors="ignore") as file:
            return file.read()

    def _split_text(self, text: str, metadata: Optional[dict] = None) -> List[Document]:
        return self.text_splitter.create_documents([text], metadatas=[metadata or {}])

    # public helper used by RAGService
    def split_text(self, text: str, metadata: Optional[dict] = None) -> List[Document]:
        """Convenience wrapper around `_split_text` used by preview logic."""
        return self._split_text(text, metadata=metadata)

    def _build_metadata(self, file_path: str, doc_type: str) -> dict:
        file_name = Path(file_path).name
        return {
            "source": file_name,
            "type": doc_type,
        }

    # utility used by preview stream
    def extract_text(self, file_path: str) -> str:
        """Return the raw text of a document without splitting."""
        ext = Path(file_path).suffix.lower()
        if ext in (".doc", ".docx"):
            doc = docx.Document(file_path)
            return "\n".join([paragraph.text for paragraph in doc.paragraphs])
        elif ext == ".pdf":
            text = ""
            if pdfplumber is not None:
                with pdfplumber.open(file_path) as pdf:
                    for page in pdf.pages:
                        text += page.extract_text() or ""
            else:
                with open(file_path, "rb") as file:
                    pdf_reader = PyPDF2.PdfReader(file)
                    for page in pdf_reader.pages:
                        text += page.extract_text() or ""
            return text
        elif ext in (".ppt", ".pptx"):
            prs = Presentation(file_path)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text
        elif ext in (".md", ".markdown"):
            return self._read_plain_text(file_path)
        else:
            raise ValueError(f"Unsupported file type: {ext}")

    def get_doc_type(self, file_path: str) -> str:
        """Map file extension to a simple type string."""
        ext = Path(file_path).suffix.lower()
        if ext in (".doc", ".docx"):
            return "word"
        elif ext == ".pdf":
            return "pdf"
        elif ext in (".ppt", ".pptx"):
            return "ppt"
        elif ext in (".md", ".markdown"):
            return "markdown"
        else:
            return "unknown"

    def process_document(self, file_path: str) -> List[Document]:
        file_extension = Path(file_path).suffix.lower()

        processors = {
            ".docx": self.process_word,
            ".doc": self.process_word,
            ".pdf": self.process_pdf,
            ".pptx": self.process_ppt,
            ".ppt": self.process_ppt,
            ".md": self.process_markdown,
            ".markdown": self.process_markdown,
        }

        processor = processors.get(file_extension)
        if not processor:
            raise ValueError(f"Unsupported file type: {file_extension}")

        return processor(file_path)
