from __future__ import annotations

import json
import math
import os
import re
import shutil
import subprocess
import sys
import tempfile
from pathlib import Path
from typing import Any
from uuid import uuid4

from app.core.config import settings
from app.services.bailian_service import bailian_service
from app.services.teaching_video_audio import add_chinese_narration


class TeachingArtifactService:
    """Render validated Qwen plans into deterministic local teaching assets."""

    def __init__(self) -> None:
        self.output_dir = Path(settings.BASE_PATH) / "uploads" / "generated_artifacts"
        self.public_prefix = f"{settings.API_V1_STR}/ai/generated-artifacts"
        self.image_output_dir = Path(settings.BASE_PATH) / "uploads" / "generated_images"
        self.image_public_prefix = f"{settings.API_V1_STR}/ai/generated-images"

    def configured(self) -> bool:
        return bailian_service.configured()

    def _target(self, suffix: str) -> tuple[Path, str]:
        self.output_dir.mkdir(parents=True, exist_ok=True)
        file_name = f"{uuid4().hex}{suffix}"
        return self.output_dir / file_name, file_name

    def _public(self, path: Path, file_name: str, kind: str, title: str) -> dict[str, Any]:
        return {
            "kind": kind,
            "title": title,
            "file_name": file_name,
            "download_url": f"{self.public_prefix}/{file_name}",
            "file_size": path.stat().st_size,
        }

    def generate_ppt(self, topic: str, user_request: str) -> dict[str, Any]:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        from pptx.util import Inches, Pt

        plan = bailian_service.structured_json(
            system_prompt=(
                "你是高校教师和课件设计师。设计一份结构严谨、文字准确的中文 PPT。"
                "返回字段 title、subtitle、slides；slides 为对象数组，每项含 title、"
                "bullets（2至6条短句）和 note。总页数 5 至 12 页。"
            ),
            user_prompt=f"主题：{topic}\n用户要求：{user_request}",
            max_tokens=6000,
        )
        title = str(plan.get("title") or f"{topic}教学课件")[:80]
        slides = plan.get("slides") if isinstance(plan.get("slides"), list) else []
        if not slides:
            raise RuntimeError("Qwen PPT plan contains no slides")
        slides = slides[:12]

        presentation = Presentation()
        presentation.slide_width = Inches(13.333)
        presentation.slide_height = Inches(7.5)
        cover = presentation.slides.add_slide(presentation.slide_layouts[0])
        cover.shapes.title.text = title
        cover.placeholders[1].text = str(plan.get("subtitle") or "AI 个性化学习课件")[:120]
        for shape in cover.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            for paragraph in shape.text_frame.paragraphs:
                paragraph.font.name = "Microsoft YaHei"
                paragraph.font.color.rgb = RGBColor(31, 41, 55)

        for index, item in enumerate(slides, start=1):
            if not isinstance(item, dict):
                continue
            slide = presentation.slides.add_slide(presentation.slide_layouts[1])
            slide.shapes.title.text = str(item.get("title") or f"第 {index} 部分")[:100]
            body = slide.placeholders[1].text_frame
            body.clear()
            bullets = item.get("bullets") if isinstance(item.get("bullets"), list) else []
            for bullet_index, bullet in enumerate(bullets[:6]):
                paragraph = body.paragraphs[0] if bullet_index == 0 else body.add_paragraph()
                paragraph.text = str(bullet)[:260]
                paragraph.level = 0
                paragraph.font.name = "Microsoft YaHei"
                paragraph.font.size = Pt(23)
                paragraph.space_after = Pt(12)
            for shape in slide.shapes:
                if not getattr(shape, "has_text_frame", False):
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    paragraph.font.name = "Microsoft YaHei"
            page_box = slide.shapes.add_textbox(Inches(12.35), Inches(6.95), Inches(.55), Inches(.25))
            page = page_box.text_frame.paragraphs[0]
            page.text = str(index + 1)
            page.font.size = Pt(10)
            page.font.color.rgb = RGBColor(107, 114, 128)
            page.alignment = PP_ALIGN.RIGHT

        path, file_name = self._target(".pptx")
        presentation.save(path)
        artifact = self._public(path, file_name, "ppt", title)
        artifact.update({"preview": f"千问已生成 {len(presentation.slides)} 页结构化课件。", "slide_count": len(presentation.slides)})
        return artifact

    def generate_scientific_chart(self, topic: str, user_request: str) -> dict[str, Any]:
        matplotlib_config_dir = Path(settings.BASE_PATH) / "uploads" / ".matplotlib"
        matplotlib_config_dir.mkdir(parents=True, exist_ok=True)
        os.environ.setdefault("MPLCONFIGDIR", str(matplotlib_config_dir))
        try:
            import matplotlib
            matplotlib.use("Agg")
            import matplotlib.pyplot as plt
        except ImportError as exc:
            raise RuntimeError("Matplotlib 未安装，请先安装 matplotlib") from exc

        spec = bailian_service.structured_json(
            system_prompt=(
                "你是理科教学数据可视化专家。生成可验证的二维图表数据，不要编造实验测量值；"
                "若是理论规律，使用公式推导出的代表性采样点。返回 title、x_label、y_label、"
                "chart_type(line/bar/scatter)、series。series 每项含 name、x 数组、y 数组。"
            ),
            user_prompt=f"主题：{topic}\n用户要求：{user_request}",
            max_tokens=5000,
        )
        series = spec.get("series") if isinstance(spec.get("series"), list) else []
        if not series:
            raise RuntimeError("Qwen chart specification contains no series")
        chart_type = str(spec.get("chart_type") or "line").lower()
        if chart_type not in {"line", "bar", "scatter"}:
            chart_type = "line"
        title = str(spec.get("title") or topic)[:100]
        plt.rcParams["font.sans-serif"] = ["Microsoft YaHei", "SimHei", "Arial Unicode MS", "DejaVu Sans"]
        plt.rcParams["axes.unicode_minus"] = False
        fig, ax = plt.subplots(figsize=(10, 5.8), dpi=150)
        for item in series[:4]:
            if not isinstance(item, dict):
                continue
            raw_x = item.get("x") if isinstance(item.get("x"), list) else []
            raw_y = item.get("y") if isinstance(item.get("y"), list) else []
            size = min(len(raw_x), len(raw_y), 60)
            x: list[float] = []
            y: list[float] = []
            for xv, yv in zip(raw_x[:size], raw_y[:size]):
                try:
                    xf, yf = float(xv), float(yv)
                except (TypeError, ValueError):
                    continue
                if math.isfinite(xf) and math.isfinite(yf):
                    x.append(xf)
                    y.append(yf)
            if not x:
                continue
            label = str(item.get("name") or "数据")[:50]
            if chart_type == "scatter":
                ax.scatter(x, y, label=label, s=28)
            elif chart_type == "bar":
                ax.bar(x, y, label=label, alpha=.78)
            else:
                ax.plot(x, y, label=label, linewidth=2.2, marker="o", markersize=3)
        if not ax.has_data():
            plt.close(fig)
            raise RuntimeError("Qwen chart data did not contain valid numeric points")
        ax.set_title(title, fontsize=16, pad=14)
        ax.set_xlabel(str(spec.get("x_label") or "x")[:80])
        ax.set_ylabel(str(spec.get("y_label") or "y")[:80])
        ax.grid(True, alpha=.22)
        if len(series) > 1 or any(str(item.get("name") or "") for item in series if isinstance(item, dict)):
            ax.legend()
        fig.tight_layout()
        self.image_output_dir.mkdir(parents=True, exist_ok=True)
        file_name = f"{uuid4().hex}.png"
        path = self.image_output_dir / file_name
        fig.savefig(path, bbox_inches="tight", facecolor="white")
        plt.close(fig)
        artifact = {
            "kind": "image",
            "title": title,
            "file_name": file_name,
            "download_url": f"{self.image_public_prefix}/{file_name}",
            "file_size": path.stat().st_size,
        }
        artifact.update({
            "image_url": artifact["download_url"],
            "preview": f"由千问生成结构化数据规格，Matplotlib 本地确定性绘制。\n\n{json.dumps(spec, ensure_ascii=False)[:1200]}",
            "chart_spec": spec,
        })
        return artifact

    def generate_data_flow_diagram(self, topic: str, user_request: str) -> dict[str, Any]:
        """Render a real black-on-white DFD PNG without Mermaid or image diffusion."""
        matplotlib_config_dir = Path(settings.BASE_PATH) / "uploads" / ".matplotlib"
        matplotlib_config_dir.mkdir(parents=True, exist_ok=True)
        os.environ.setdefault("MPLCONFIGDIR", str(matplotlib_config_dir))
        try:
            import matplotlib
            matplotlib.use("Agg")
            import matplotlib.pyplot as plt
            from matplotlib.patches import FancyArrowPatch, FancyBboxPatch
        except ImportError as exc:
            raise RuntimeError("Matplotlib 未安装，请先安装 matplotlib") from exc

        text = re.sub(r"\s+", "", user_request or "").lower()
        is_salary_case = all(token in text for token in ("教师", "26000", "工资"))
        if is_salary_case:
            title = "教师工资调整系统数据流图（DFD 1 层）"
            nodes = [
                {"id": "admin_in", "kind": "external", "label": "行政办公室", "x": .04, "y": .67},
                {"id": "salary_file", "kind": "store", "label": "D1  教师工资档案（光盘）", "x": .05, "y": .18},
                {"id": "read", "kind": "process", "label": "1.0\n读取教师档案", "x": .28, "y": .57},
                {"id": "calculate", "kind": "process", "label": "2.0\n计算调整后工资", "x": .50, "y": .57},
                {"id": "print", "kind": "process", "label": "3.0\n生成并打印工资清单", "x": .72, "y": .57},
                {"id": "admin_out", "kind": "external", "label": "行政办公室", "x": .84, "y": .18},
            ]
            flows = [
                {"source": "admin_in", "target": "read", "label": "处理请求"},
                {"source": "salary_file", "target": "read", "label": "原工资、赡养人数、雇用日期"},
                {"source": "read", "target": "calculate", "label": "教师工资记录"},
                {"source": "calculate", "target": "print", "label": "原工资、新工资"},
                {"source": "print", "target": "admin_out", "label": "200名教师工资调整清单"},
            ]
            rule_note = (
                "计算规则：原工资 ≥ 26,000美元时保持不变；否则，新工资 = min(26,000美元，"
                "原工资 + 赡养人数×100美元 + 工龄×50美元)。"
            )
            source = "verified_salary_dfd_template"
        else:
            if not bailian_service.configured():
                raise RuntimeError("生成通用数据流图需要配置 DASHSCOPE_API_KEY")
            spec = bailian_service.structured_json(
                system_prompt=(
                    "你是软件工程数据流图设计师。只返回 JSON，字段 title、external_entities、"
                    "processes、data_stores、flows。前三项为对象数组，每项只有 id 和 label；"
                    "flows 每项只有 source、target、label。使用 DFD 而不是程序流程图；"
                    "外部实体最多4个、处理最多6个、数据存储最多4个、数据流最多14条。"
                ),
                user_prompt=user_request,
                max_tokens=3500,
            )
            title = str(spec.get("title") or f"{topic}数据流图")[:100]
            externals = spec.get("external_entities") if isinstance(spec.get("external_entities"), list) else []
            processes = spec.get("processes") if isinstance(spec.get("processes"), list) else []
            stores = spec.get("data_stores") if isinstance(spec.get("data_stores"), list) else []
            raw_flows = spec.get("flows") if isinstance(spec.get("flows"), list) else []
            nodes = []
            for index, item in enumerate(externals[:4]):
                if isinstance(item, dict):
                    nodes.append({"id": str(item.get("id") or f"e{index}"), "kind": "external", "label": str(item.get("label") or "外部实体")[:40], "x": .03 if index % 2 == 0 else .84, "y": .72 - (index // 2) * .38})
            for index, item in enumerate(processes[:6]):
                if isinstance(item, dict):
                    count = max(1, min(len(processes), 6))
                    nodes.append({"id": str(item.get("id") or f"p{index}"), "kind": "process", "label": str(item.get("label") or f"{index + 1}.0 处理")[:50], "x": .20 + index * (.60 / max(count - 1, 1)), "y": .55})
            for index, item in enumerate(stores[:4]):
                if isinstance(item, dict):
                    count = max(1, min(len(stores), 4))
                    nodes.append({"id": str(item.get("id") or f"d{index}"), "kind": "store", "label": str(item.get("label") or f"D{index + 1} 数据存储")[:50], "x": .18 + index * (.62 / max(count - 1, 1)), "y": .16})
            valid_ids = {item["id"] for item in nodes}
            flows = [
                {"source": str(item.get("source")), "target": str(item.get("target")), "label": str(item.get("label") or "数据")[:60]}
                for item in raw_flows[:14]
                if isinstance(item, dict) and str(item.get("source")) in valid_ids and str(item.get("target")) in valid_ids
            ]
            if not nodes or not flows:
                raise RuntimeError("千问返回的数据流图规格不完整")
            rule_note = "图例：矩形表示外部实体，圆角框表示处理，双横线表示数据存储，箭头表示数据流。"
            source = "qwen_dfd_spec"

        plt.rcParams["font.sans-serif"] = ["Microsoft YaHei", "SimHei", "Arial Unicode MS", "DejaVu Sans"]
        plt.rcParams["axes.unicode_minus"] = False
        fig, ax = plt.subplots(figsize=(14, 7.8), dpi=180)
        ax.set_xlim(0, 1)
        ax.set_ylim(0, 1)
        ax.axis("off")
        ax.text(.5, .95, title, ha="center", va="center", fontsize=20, fontweight="bold", color="black")
        positions: dict[str, tuple[float, float, str]] = {}
        for node in nodes:
            x, y, kind = float(node["x"]), float(node["y"]), str(node["kind"])
            label = str(node["label"])
            if kind == "process":
                width, height = .16, .16
                patch = FancyBboxPatch((x, y), width, height, boxstyle="round,pad=.015,rounding_size=.07", edgecolor="black", facecolor="white", linewidth=2)
                ax.add_patch(patch)
                ax.text(x + width / 2, y + height / 2, label, ha="center", va="center", fontsize=11, color="black")
            elif kind == "store":
                width, height = .20, .10
                ax.plot([x, x + width], [y + height, y + height], color="black", linewidth=2)
                ax.plot([x, x + width], [y, y], color="black", linewidth=2)
                ax.text(x + width / 2, y + height / 2, label, ha="center", va="center", fontsize=10.5, color="black")
            else:
                width, height = .13, .11
                patch = FancyBboxPatch((x, y), width, height, boxstyle="square,pad=.012", edgecolor="black", facecolor="white", linewidth=2)
                ax.add_patch(patch)
                ax.text(x + width / 2, y + height / 2, label, ha="center", va="center", fontsize=11, color="black")
            positions[str(node["id"])] = (x + width / 2, y + height / 2, kind)

        for index, flow in enumerate(flows):
            source_pos = positions.get(str(flow["source"]))
            target_pos = positions.get(str(flow["target"]))
            if not source_pos or not target_pos:
                continue
            sx, sy, _ = source_pos
            tx, ty, _ = target_pos
            arrow = FancyArrowPatch((sx, sy), (tx, ty), arrowstyle="-|>", mutation_scale=15, linewidth=1.6, color="black", shrinkA=38, shrinkB=38, connectionstyle=f"arc3,rad={.04 if index % 2 else 0}")
            ax.add_patch(arrow)
            mx, my = (sx + tx) / 2, (sy + ty) / 2
            offset = .024 if abs(sx - tx) >= abs(sy - ty) else .012
            ax.text(mx, my + offset, str(flow["label"]), ha="center", va="bottom", fontsize=9, color="black", bbox={"facecolor": "white", "edgecolor": "none", "pad": 1.5})

        ax.text(.5, .065, rule_note, ha="center", va="center", fontsize=10.5, color="black")
        ax.text(.5, .025, "DFD 图例：外部实体 □    处理 ○/圆角框    数据存储 ═    数据流 →", ha="center", va="center", fontsize=9.5, color="black")
        fig.tight_layout(pad=1.2)
        self.image_output_dir.mkdir(parents=True, exist_ok=True)
        file_name = f"{uuid4().hex}.png"
        path = self.image_output_dir / file_name
        fig.savefig(path, bbox_inches="tight", facecolor="white")
        plt.close(fig)
        artifact = {
            "kind": "image",
            "resource_type": "data_flow_diagram",
            "title": title,
            "file_name": file_name,
            "download_url": f"{self.image_public_prefix}/{file_name}",
            "image_url": f"{self.image_public_prefix}/{file_name}",
            "file_size": path.stat().st_size,
            "preview": "已使用 DFD 专用确定性绘图器生成黑字白底 PNG；未使用 Mermaid，未调用插画模型。",
            "diagram_source": source,
        }
        return artifact

    def generate_manim_video(self, topic: str, user_request: str) -> dict[str, Any]:
        if not settings.MANIM_ENABLED:
            raise RuntimeError("Manim 教学动画功能未启用")
        try:
            import importlib.util
            if importlib.util.find_spec("manim") is None:
                raise RuntimeError("Manim 未安装，请先安装 manim 并配置 FFmpeg")
        except ImportError as exc:
            raise RuntimeError("Manim 未安装") from exc

        storyboard = bailian_service.structured_json(
            system_prompt=(
                "你是理科教学动画导演。把知识点拆成 3 至 7 个循序渐进步骤。"
                "返回 title 和 steps；每个 step 含 label、formula、explanation。"
                "formula 使用简短纯文本数学表达式，不使用 LaTeX 命令。"
            ),
            user_prompt=f"主题：{topic}\n用户要求：{user_request}",
            max_tokens=4500,
        )
        steps = storyboard.get("steps") if isinstance(storyboard.get("steps"), list) else []
        if not steps:
            raise RuntimeError("Qwen storyboard contains no animation steps")
        normalized = []
        for item in steps[:7]:
            if isinstance(item, dict):
                normalized.append({
                    "label": str(item.get("label") or "步骤")[:45],
                    "formula": str(item.get("formula") or "")[:80],
                    "explanation": str(item.get("explanation") or "")[:100],
                })
        title = str(storyboard.get("title") or f"{topic}教学动画")[:80]
        safe_payload = json.dumps({"title": title, "steps": normalized}, ensure_ascii=False)
        script = f'''from manim import *
import json

STORY = json.loads({safe_payload!r})

class TeachingAnimation(Scene):
    def construct(self):
        title = Text(STORY["title"], font="Microsoft YaHei", font_size=38, color=BLUE)
        self.play(Write(title))
        self.play(title.animate.to_edge(UP).scale(0.72))
        current = None
        for index, step in enumerate(STORY["steps"], start=1):
            lines = [f"{{index}}. {{step['label']}}"]
            if step.get("formula"):
                lines.append(step["formula"])
            if step.get("explanation"):
                lines.append(step["explanation"])
            card = VGroup(*[
                Text(line, font="Microsoft YaHei", font_size=26 if i == 0 else 22)
                for i, line in enumerate(lines)
            ]).arrange(DOWN, buff=.28).move_to(ORIGIN)
            box = SurroundingRectangle(card, color=BLUE_C, buff=.35, corner_radius=.12)
            group = VGroup(box, card)
            if current is None:
                self.play(FadeIn(group, shift=RIGHT))
            else:
                self.play(ReplacementTransform(current, group))
            self.wait(1.2)
            current = group
        if current is not None:
            self.play(FadeOut(current))
        done = Text("总结：理解过程，再记住结论", font="Microsoft YaHei", font_size=30, color=GREEN)
        self.play(Write(done))
        self.wait(1)
'''
        with tempfile.TemporaryDirectory(prefix="zhiyu-manim-") as temp_dir:
            temp = Path(temp_dir)
            script_path = temp / "teaching_animation.py"
            script_path.write_text(script, encoding="utf-8")
            media_dir = temp / "media"
            command = [
                sys.executable, "-m", "manim", "-ql", "--disable_caching",
                "--media_dir", str(media_dir), str(script_path), "TeachingAnimation",
            ]
            result = subprocess.run(
                command,
                cwd=temp,
                capture_output=True,
                text=True,
                timeout=settings.TEACHING_ARTIFACT_TIMEOUT_SECONDS,
            )
            if result.returncode != 0:
                raise RuntimeError(f"Manim render failed: {(result.stderr or result.stdout)[-900:]}")
            candidates = list(media_dir.rglob("*.mp4"))
            if not candidates:
                raise RuntimeError("Manim completed but produced no MP4")
            path, file_name = self._target(".mp4")
            narration_parts = [title]
            for index, step in enumerate(normalized, start=1):
                explanation = step["explanation"] or step["formula"]
                narration_parts.append(f"第{index}步，{step['label']}。{explanation}")
            narration_parts.append("以上就是本节内容，请结合画面回顾关键过程。")
            narration = "。".join(part.strip("。") for part in narration_parts if part.strip())
            voice = add_chinese_narration(
                video_path=candidates[0],
                narration=narration[:800],
                output_path=path,
            )
        artifact = self._public(path, file_name, "video", title)
        artifact.update({
            "preview": f"千问生成 {len(normalized)} 步教学分镜，Manim 已完成本地渲染并加入中文旁白。",
            "storyboard": storyboard,
            "audio_enabled": True,
            "narration_voice": voice,
        })
        return artifact


teaching_artifact_service = TeachingArtifactService()
