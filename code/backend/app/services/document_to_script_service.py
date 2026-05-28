from __future__ import annotations

import re
from dataclasses import asdict, dataclass
from pathlib import Path

import PyPDF2
from pptx import Presentation

from app.services.digital_human_assets import DEFAULT_GESTURES
from app.services.document_processor import DocumentProcessor


@dataclass
class ScriptSegment:
    index: int
    title: str
    narration: str
    gesture_id: str
    duration_seconds: float
    source_ref: str


@dataclass
class DigitalHumanScript:
    title: str
    narration: str
    segments: list[ScriptSegment]
    gesture_timeline: list[dict]

    def to_dict(self) -> dict:
        return {
            "title": self.title,
            "narration": self.narration,
            "segments": [asdict(item) for item in self.segments],
            "gesture_timeline": self.gesture_timeline,
        }


class DocumentToScriptService:
    max_segments = 8
    max_chars_per_segment = 480

    def build_text_script(self, text: str, title: str | None = None) -> DigitalHumanScript:
        body = " ".join((text or "").split())
        if not body:
            raise RuntimeError("文本生成视频任务缺少脚本文本")
        chunks = self._chunk_plain_text(body)
        return self._build_script(chunks, title or "数字人讲解", source_prefix="文本")

    def build_document_script(self, source_path: str, title: str | None = None) -> DigitalHumanScript:
        path = Path(source_path)
        page_chunks = self._extract_page_chunks(path)
        if not page_chunks:
            processor = DocumentProcessor()
            fallback = " ".join((processor.extract_text(str(path)) or "").split())
            page_chunks = self._chunk_plain_text(fallback)
        if not page_chunks:
            raise RuntimeError("上传的课件未提取到可用文本")
        return self._build_script(page_chunks, title or path.stem, source_prefix="页面")

    def _build_script(
        self,
        chunks: list[tuple[str, str]],
        title: str,
        *,
        source_prefix: str,
    ) -> DigitalHumanScript:
        segments: list[ScriptSegment] = []
        elapsed = 0.0
        timeline: list[dict] = []

        for index, (source_ref, raw_text) in enumerate(chunks[: self.max_segments], start=1):
            text = self._clean_segment_text(raw_text)
            if not text:
                continue
            gesture_id = self._select_gesture(text, index)
            narration = self._segment_narration(index, text, source_ref, source_prefix)
            duration = self._estimate_duration(narration)
            segments.append(
                ScriptSegment(
                    index=index,
                    title=f"{source_prefix}{index}: {self._segment_title(text)}",
                    narration=narration,
                    gesture_id=gesture_id,
                    duration_seconds=duration,
                    source_ref=source_ref,
                )
            )
            timeline.append(
                {
                    "time": round(elapsed, 2),
                    "gesture_id": gesture_id,
                    "label": self._gesture_label(gesture_id),
                    "segment_index": index,
                }
            )
            elapsed += duration

        if not segments:
            raise RuntimeError("未生成可用数字人讲解脚本")

        opening = f"同学你好，下面我们用几分钟讲清楚《{title}》的核心内容。"
        closing = "最后请你回到 AI 伴学里完成一道变式练习，我会根据作答继续调整学习建议。"
        narration = "\n".join([opening, *[item.narration for item in segments], closing])
        return DigitalHumanScript(
            title=title,
            narration=narration,
            segments=segments,
            gesture_timeline=timeline,
        )

    def _extract_page_chunks(self, path: Path) -> list[tuple[str, str]]:
        ext = path.suffix.lower()
        if ext in (".ppt", ".pptx"):
            return self._extract_ppt_chunks(path)
        if ext == ".pdf":
            return self._extract_pdf_chunks(path)
        return []

    def _extract_ppt_chunks(self, path: Path) -> list[tuple[str, str]]:
        prs = Presentation(str(path))
        chunks: list[tuple[str, str]] = []
        for index, slide in enumerate(prs.slides, start=1):
            texts: list[str] = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    value = str(shape.text or "").strip()
                    if value:
                        texts.append(value)
            merged = " ".join(" ".join(texts).split())
            if merged:
                chunks.append((f"第 {index} 页", merged))
        return chunks

    def _extract_pdf_chunks(self, path: Path) -> list[tuple[str, str]]:
        chunks: list[tuple[str, str]] = []
        with open(path, "rb") as file:
            reader = PyPDF2.PdfReader(file)
            for index, page in enumerate(reader.pages, start=1):
                text = " ".join((page.extract_text() or "").split())
                if text:
                    chunks.append((f"第 {index} 页", text))
        return chunks

    def _chunk_plain_text(self, text: str) -> list[tuple[str, str]]:
        clean = self._clean_segment_text(text)
        if not clean:
            return []
        sentences = re.split(r"(?<=[。！？!?；;])\s*", clean)
        chunks: list[str] = []
        current = ""
        for sentence in sentences:
            if not sentence:
                continue
            if current and len(current) + len(sentence) > self.max_chars_per_segment:
                chunks.append(current)
                current = sentence
            else:
                current += sentence
        if current:
            chunks.append(current)
        return [(f"第 {index} 段", chunk) for index, chunk in enumerate(chunks, start=1)]

    def _segment_narration(self, index: int, text: str, source_ref: str, source_prefix: str) -> str:
        trimmed = text[: self.max_chars_per_segment]
        lead = "首先" if index == 1 else "接下来"
        if source_prefix == "页面":
            return f"{lead}看{source_ref}。{trimmed}"
        return f"{lead}我们看第 {index} 个要点。{trimmed}"

    @staticmethod
    def _clean_segment_text(text: str) -> str:
        return re.sub(r"\s+", " ", (text or "").strip())

    @staticmethod
    def _segment_title(text: str) -> str:
        stripped = re.sub(r"[#*`>\\-]+", "", text).strip()
        return stripped[:18] or "核心内容"

    @staticmethod
    def _estimate_duration(text: str) -> float:
        return max(5.0, min(18.0, len(text) / 9.0))

    @staticmethod
    def _gesture_label(gesture_id: str) -> str:
        for item in DEFAULT_GESTURES:
            if item.id == gesture_id:
                return item.label
        return "自然讲解"

    @staticmethod
    def _select_gesture(text: str, index: int) -> str:
        if re.search(r"对比|区别|相同|不同|一方面|另一方面", text):
            return "compare_two_sides"
        if re.search(r"重点|关键|注意|必须|核心", text):
            return "emphasis_one_hand"
        if re.search(r"左侧|第一步|首先|目录", text):
            return "point_left"
        if re.search(r"右侧|公式|步骤|流程|图表", text):
            return "point_right"
        if re.search(r"总结|归纳|最后|因此|所以", text):
            return "nod_summary"
        if re.search(r"练习|作答|批改|继续|掌握", text):
            return "encourage_forward"
        if index % 3 == 2:
            return "explain_open"
        return "idle"


document_to_script_service = DocumentToScriptService()
