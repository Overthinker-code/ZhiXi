from pathlib import Path
from types import SimpleNamespace
from uuid import uuid4

import pytest
from docx import Document
from pptx import Presentation
from sqlalchemy import create_engine
from sqlmodel import SQLModel, Session

from app import models
from app.api.routes import resources
from fastapi.responses import FileResponse
from app.db.base_class import Base
from app.services.resource_preview_service import ResourcePreviewError, ResourcePreviewService


def _session() -> Session:
    engine = create_engine("sqlite:///:memory:")
    SQLModel.metadata.create_all(engine)
    Base.metadata.create_all(engine)
    return Session(engine)


def _resource(owner_id, *, name: str, source: str | None = None) -> models.Resource:
    return models.Resource(
        title="安全预览资料",
        type="docx",
        subject="数据库",
        file_name=name,
        file_path=name,
        content_type="application/octet-stream",
        source=source,
        uploader_id=owner_id,
    )


def _owner(user_id, *, is_superuser: bool = False):
    return SimpleNamespace(id=user_id, is_superuser=is_superuser)


def test_pdf_preview_is_inline_private_and_records_evidence_once(tmp_path: Path, monkeypatch) -> None:
    db = _session()
    owner_id = uuid4()
    resource = _resource(owner_id, name="handout.pdf")
    db.add(resource)
    db.commit()
    db.refresh(resource)
    (tmp_path / "handout.pdf").write_bytes(b"%PDF-1.4\npreview")
    monkeypatch.setattr(resources, "UPLOAD_DIR", str(tmp_path))
    monkeypatch.setattr(
        Path,
        "read_bytes",
        lambda _path: (_ for _ in ()).throw(AssertionError("binary preview must stream")),
    )
    evidence: list[str] = []
    monkeypatch.setattr(resources, "_record_resource_event", lambda *_args, **kwargs: evidence.append(kwargs["event_type"]))

    response = resources.preview_resource(db=db, current_user=_owner(owner_id), resource_id=resource.id)

    assert response.status_code == 200
    assert isinstance(response, FileResponse)
    assert response.media_type == "application/pdf"
    assert response.headers["content-disposition"].startswith("inline")
    assert response.headers["cache-control"] == "private, no-store"
    assert response.headers["x-content-type-options"] == "nosniff"
    assert "default-src 'none'" in response.headers["content-security-policy"]
    assert response.path.endswith("handout.pdf")
    assert evidence == ["resource_previewed"]
    assert "file_path" not in resources._resource_payload(resource)


def test_preview_headers_support_chinese_filename_with_rfc5987_encoding() -> None:
    headers = resources._preview_headers("事务讲义（学生版）.docx")

    assert 'filename="resource.docx"' in headers["Content-Disposition"]
    assert "filename*=UTF-8''%E4%BA%8B%E5%8A%A1" in headers["Content-Disposition"]
    assert "\u4e8b\u52a1" not in headers["Content-Disposition"]


def test_preview_headers_keep_converted_html_inert_without_blocking_blob_iframe() -> None:
    headers = resources._preview_headers("slides.pptx")
    csp = headers["Content-Security-Policy"]

    assert headers["Cache-Control"] == "private, no-store"
    assert headers["X-Content-Type-Options"] == "nosniff"
    assert "default-src 'none'" in csp
    assert "style-src 'unsafe-inline'" in csp
    assert "img-src data: blob:" in csp
    assert "base-uri 'none'" in csp
    assert "form-action 'none'" in csp
    assert "frame-ancestors" not in csp


def test_docx_preview_escapes_content_and_keeps_structural_html(tmp_path: Path, monkeypatch) -> None:
    db = _session()
    owner_id = uuid4()
    resource = _resource(owner_id, name="notes.docx")
    db.add(resource)
    db.commit()
    db.refresh(resource)
    document = Document()
    document.add_heading("<script>alert(1)</script>", level=1)
    document.add_paragraph("第一段")
    document.add_paragraph("列表项", style="List Bullet")
    table = document.add_table(rows=1, cols=2)
    table.cell(0, 0).text = "字段"
    table.cell(0, 1).text = "值"
    document.save(tmp_path / "notes.docx")
    monkeypatch.setattr(resources, "UPLOAD_DIR", str(tmp_path))
    monkeypatch.setattr(resources, "_record_resource_event", lambda *_args, **_kwargs: None)

    response = resources.preview_resource(db=db, current_user=_owner(owner_id), resource_id=resource.id)
    body = response.body.decode("utf-8")

    assert response.status_code == 200
    assert response.media_type.startswith("text/html")
    assert "&lt;script&gt;alert(1)&lt;/script&gt;" in body
    assert "<script>alert" not in body
    assert "<h1>" in body and "<ul>" in body and "<table>" in body
    assert "Content-Security-Policy" in body
    assert "*,*::before,*::after{box-sizing:border-box}" in body
    assert "width:min(210mm,100%)" in body
    assert "min-height:297mm" in body
    assert ".preview-presentation" in body


def test_preview_hides_private_resources_and_handles_missing_or_unsupported_files(
    tmp_path: Path, monkeypatch
) -> None:
    db = _session()
    owner_id, other_id = uuid4(), uuid4()
    private = _resource(owner_id, name="private.pdf")
    missing = _resource(owner_id, name="missing.pdf")
    legacy = _resource(owner_id, name="legacy.doc")
    db.add(private)
    db.add(missing)
    db.add(legacy)
    db.commit()
    for resource in (private, missing, legacy):
        db.refresh(resource)
    (tmp_path / "private.pdf").write_bytes(b"%PDF-1.4")
    (tmp_path / "legacy.doc").write_bytes(b"legacy")
    monkeypatch.setattr(resources, "UPLOAD_DIR", str(tmp_path))
    monkeypatch.setattr(resources, "_record_resource_event", lambda *_args, **_kwargs: None)

    with pytest.raises(Exception) as hidden:
        resources.preview_resource(db=db, current_user=_owner(other_id), resource_id=private.id)
    assert getattr(hidden.value, "status_code", None) == 404
    with pytest.raises(Exception) as absent:
        resources.preview_resource(db=db, current_user=_owner(owner_id), resource_id=missing.id)
    assert getattr(absent.value, "status_code", None) == 404
    with pytest.raises(Exception) as unsupported:
        resources.preview_resource(db=db, current_user=_owner(owner_id), resource_id=legacy.id)
    assert getattr(unsupported.value, "status_code", None) == 415


def test_non_owner_cannot_list_or_access_ordinary_uploads(tmp_path: Path, monkeypatch) -> None:
    db = _session()
    owner_id, other_id = uuid4(), uuid4()
    own_resource = _resource(other_id, name="mine.pdf")
    private_resource = _resource(owner_id, name="ordinary-upload.pdf")
    db.add(own_resource)
    db.add(private_resource)
    db.commit()
    for resource in (own_resource, private_resource):
        db.refresh(resource)
    (tmp_path / "ordinary-upload.pdf").write_bytes(b"%PDF-1.4")
    monkeypatch.setattr(resources, "UPLOAD_DIR", str(tmp_path))
    monkeypatch.setattr(resources, "_record_resource_event", lambda *_args, **_kwargs: None)

    listed = resources.read_resources(db=db, current_user=_owner(other_id))
    assert [item.id for item in listed.data] == [own_resource.id]

    protected_calls = (
        lambda: resources.read_resource(db=db, current_user=_owner(other_id), resource_id=private_resource.id),
        lambda: resources.download_resource(db=db, current_user=_owner(other_id), resource_id=private_resource.id),
        lambda: resources.preview_resource(db=db, current_user=_owner(other_id), resource_id=private_resource.id),
        lambda: resources.set_resource_favorite(
            db=db, current_user=_owner(other_id), resource_id=private_resource.id,
            update=resources.FavoriteUpdate(favorite=True),
        ),
        lambda: resources.set_resource_config(
            db=db, current_user=_owner(other_id), resource_id=private_resource.id,
            update=resources.ResourceConfigUpdate(is_top=True),
        ),
        lambda: resources.remove_resource_from_library(
            db=db, current_user=_owner(other_id), resource_id=private_resource.id,
        ),
    )
    for call in protected_calls:
        with pytest.raises(Exception) as denied:
            call()
        assert getattr(denied.value, "status_code", None) == 404


def test_superuser_can_list_and_read_another_users_resource() -> None:
    db = _session()
    owner_id, admin_id = uuid4(), uuid4()
    resource = _resource(owner_id, name="ordinary-upload.pdf")
    db.add(resource)
    db.commit()
    db.refresh(resource)

    listed = resources.read_resources(db=db, current_user=_owner(admin_id, is_superuser=True))
    assert [item.id for item in listed.data] == [resource.id]
    read = resources.read_resource(
        db=db, current_user=_owner(admin_id, is_superuser=True), resource_id=resource.id
    )
    assert read["id"] == resource.id


def test_preview_maps_file_race_oserror_to_not_found(tmp_path: Path, monkeypatch) -> None:
    db = _session()
    owner_id = uuid4()
    resource = _resource(owner_id, name="racing.pdf")
    db.add(resource)
    db.commit()
    db.refresh(resource)
    (tmp_path / "racing.pdf").write_bytes(b"%PDF-1.4")
    monkeypatch.setattr(resources, "UPLOAD_DIR", str(tmp_path))
    monkeypatch.setattr(
        resources.resource_preview_service,
        "prepare",
        lambda _path: (_ for _ in ()).throw(OSError("file vanished")),
    )

    with pytest.raises(Exception) as unavailable:
        resources.preview_resource(db=db, current_user=_owner(owner_id), resource_id=resource.id)
    assert getattr(unavailable.value, "status_code", None) == 404
    assert getattr(unavailable.value, "detail", None) == "资源文件不存在或已被删除"


def test_text_mmd_and_pptx_previews_are_escaped_and_cached_until_file_changes(tmp_path: Path) -> None:
    service = ResourcePreviewService()
    text_path = tmp_path / "diagram.mmd"
    text_path.write_text("graph TD\nA[<script>] --> B", encoding="utf-8")

    first = service.prepare(text_path)
    second = service.prepare(text_path)
    assert first.mode == "mermaid"
    assert first is second
    assert "&lt;script&gt;" in (first.content or "")

    text_path.write_text("graph TD\nChanged --> B", encoding="utf-8")
    changed = service.prepare(text_path)
    assert changed is not first
    assert "Changed" in (changed.content or "")

    presentation = Presentation()
    presentation.slides.add_slide(presentation.slide_layouts[1]).shapes.title.text = "PPTX 预览"
    pptx_path = tmp_path / "slides.pptx"
    presentation.save(pptx_path)
    assert service.prepare(pptx_path).mode == "presentation"


def test_preview_rejects_oversized_and_invalid_office_files(tmp_path: Path) -> None:
    service = ResourcePreviewService()
    oversized = tmp_path / "large.txt"
    oversized.write_bytes(b"x" * (service._MAX_FILE_BYTES + 1))
    with pytest.raises(ResourcePreviewError, match="文件过大"):
        service.prepare(oversized)

    invalid = tmp_path / "broken.docx"
    invalid.write_bytes(b"not an office zip")
    with pytest.raises(ResourcePreviewError, match="无法安全解析"):
        service.prepare(invalid)
